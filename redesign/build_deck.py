from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

# ---------- palette ----------
NAVY    = RGBColor(0x15, 0x22, 0x4A)
NAVY_BG = RGBColor(0x10, 0x1A, 0x38)
PANEL   = RGBColor(0x1B, 0x2A, 0x57)
INK     = RGBColor(0x2E, 0x35, 0x48)
GRAY    = RGBColor(0x5A, 0x62, 0x73)
ICE     = RGBColor(0xC9, 0xD6, 0xF0)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
TEAL    = RGBColor(0x0E, 0x8F, 0x8B)
TEAL_D  = RGBColor(0x0B, 0x6E, 0x6B)
TEAL_T  = RGBColor(0xE4, 0xF2, 0xF1)
AMBER   = RGBColor(0xE3, 0xA9, 0x1C)
AMBER_D = RGBColor(0x9E, 0x70, 0x08)
AMBER_T = RGBColor(0xFB, 0xF2, 0xDC)
NAVY_T  = RGBColor(0xEE, 0xF1, 0xF8)
LINE    = RGBColor(0xD8, 0xDD, 0xE8)

HEAD, BODY = "Georgia", "Calibri"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s

def text(s, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP, wrap=True):
    """paras: list of dicts {runs:[(txt, opts)], align, space_before, space_after, line}"""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, p in enumerate(paras):
        par = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        par.alignment = p.get("align", PP_ALIGN.LEFT)
        if "space_before" in p: par.space_before = Pt(p["space_before"])
        if "space_after"  in p: par.space_after  = Pt(p["space_after"])
        if "line" in p: par.line_spacing = p["line"]
        for txt, o in p["runs"]:
            r = par.add_run(); r.text = txt
            f = r.font
            f.name = o.get("font", BODY)
            f.size = Pt(o.get("size", 14))
            f.bold = o.get("bold", False)
            f.italic = o.get("italic", False)
            f.color.rgb = o.get("color", INK)
    return tb

def card(s, x, y, w, h, fill, line=None, radius=0.055):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                            Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = radius
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = Pt(1)
    else: sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def circle(s, x, y, d, fill, txt=None, tcolor=WHITE, tsize=15, font=HEAD, bold=True, italic=False):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background(); sh.shadow.inherit = False
    if txt:
        tf = sh.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = txt
        r.font.name = font; r.font.size = Pt(tsize); r.font.bold = bold
        r.font.italic = italic; r.font.color.rgb = tcolor
    return sh

def img_fit(s, path, x, y, w, h):
    """place image centered inside box, preserving aspect"""
    iw, ih = Image.open(path).size
    ar, box = iw / ih, w / h
    if ar > box: dw, dh = w, w / ar
    else:        dw, dh = h * ar, h
    return s.shapes.add_picture(path, Inches(x + (w - dw) / 2),
                                Inches(y + (h - dh) / 2), Inches(dw), Inches(dh))

def kicker_title(s, kicker, title, kcolor=TEAL, tcolor=NAVY, x=0.9, y=0.5, w=11.5):
    text(s, x, y, w, 0.3, [{"runs": [(kicker.upper(), dict(font=BODY, size=12, bold=True, color=kcolor))]}])
    text(s, x, y + 0.32, w, 0.75, [{"runs": [(title, dict(font=HEAD, size=30, bold=True, color=tcolor))]}])

def pageno(s, n, dark=False):
    text(s, 12.55, 7.08, 0.5, 0.3,
         [{"runs": [(str(n), dict(font=BODY, size=9, color=ICE if dark else GRAY))],
           "align": PP_ALIGN.RIGHT}])

A = "assets/"

# ================= SLIDE 1 — TITLE (dark) =================
s = slide(NAVY_BG)
text(s, 0.9, 0.85, 7.0, 0.3, [{"runs": [("B.SC. RESEARCH PROJECT  ·  SUMMARY REPORT",
     dict(font=BODY, size=12, bold=True, color=RGBColor(0x4F, 0xD1, 0xCC)))]}])
text(s, 0.9, 1.3, 6.9, 2.6, [{"runs": [("Risk-Adjusted Performance Analysis Based on Data",
     dict(font=HEAD, size=40, bold=True, color=WHITE))], "line": 1.05}])
text(s, 0.9, 3.75, 6.4, 1.1, [{"runs": [
    ("Five classic ratios, one new metric — ", dict(size=15, color=ICE)),
    ("diff × r × b", dict(font=HEAD, size=15, italic=True, color=RGBColor(0x4F, 0xD1, 0xCC))),
    (" — tested across 100 ETFs and 9 market regimes.", dict(size=15, color=ICE))], "line": 1.25}])
img_fit(s, A + "hero.png", 7.55, 1.45, 5.15, 3.4)
text(s, 0.9, 5.95, 11.5, 0.35, [{"runs": [("Nehoray Sade   ·   Aviv Asraf",
     dict(font=BODY, size=15, bold=True, color=WHITE))]}])
text(s, 0.9, 6.38, 11.5, 0.6, [{"runs": [
    ("Supervised by Dr. Dvir Ross   ·   Department of Computer Science, Sami Shamoon College of Engineering, Be'er Sheva   ·   June 2, 2026",
     dict(size=11.5, color=ICE))]}])

# ================= SLIDE 2 — ABSTRACT =================
s = slide()
kicker_title(s, "Abstract", "One new metric against five classics")
text(s, 0.9, 1.75, 6.6, 1.15, [{"runs": [
    ("The project examines the ", dict(size=14.5)),
    ("Sharpe, Treynor, Calmar, Sortino", dict(size=14.5, bold=True)),
    (" and ", dict(size=14.5)),
    ("Information", dict(size=14.5, bold=True)),
    (" ratios, and proposes a new metric derived from an exponential regression on the asset price.",
     dict(size=14.5))], "line": 1.2}])
card(s, 0.9, 3.0, 6.6, 1.25, TEAL_T)
img_fit(s, A + "f_score.png", 1.4, 3.2, 5.6, 0.85)
text(s, 0.9, 4.6, 6.6, 2.3, [
    {"runs": [("Across 100 ETFs and 9 market regimes, the new metric acts as a ", dict(size=14.5)),
              ("mean-reversion signal", dict(size=14.5, bold=True, color=TEAL_D)),
              (", while the classic ratios behave as ", dict(size=14.5)),
              ("momentum signals", dict(size=14.5, bold=True, color=AMBER_D)),
              (".", dict(size=14.5))], "line": 1.25, "space_after": 8},
    {"runs": [("Combining them according to the market regime outperforms each one individually.",
               dict(size=14.5, bold=True, color=NAVY))], "line": 1.25}])
comp = [("diff", "Price deviation", "Distance of the current price from the trend line"),
        ("r",    "Fit quality",     "Correlation coefficient — quality of fit to the trend"),
        ("b",    "Growth slope",    "Daily growth rate of the fitted trend")]
cy = 1.75
for sym, lab, desc in comp:
    card(s, 8.0, cy, 4.45, 1.52, NAVY_T)
    circle(s, 8.3, cy + 0.42, 0.68, TEAL, sym, tsize=17 if len(sym) > 1 else 20, italic=True)
    text(s, 9.25, cy + 0.26, 3.0, 0.35, [{"runs": [(lab, dict(font=HEAD, size=15, bold=True, color=NAVY))]}])
    text(s, 9.25, cy + 0.64, 3.05, 0.8, [{"runs": [(desc, dict(size=11.5, color=GRAY))], "line": 1.1}])
    cy += 1.72
pageno(s, 2)

# ================= SLIDE 3 — FIVE CLASSIC RATIOS =================
s = slide()
kicker_title(s, "The momentum family", "The five classic ratios", kcolor=AMBER_D)
text(s, 0.9, 1.62, 11.5, 0.4, [{"runs": [
    ("Each ratio highlights a different aspect of risk — and all five reward strong recent risk-adjusted returns.",
     dict(size=14, color=GRAY))]}])
ratios = [("Sharpe", "1966", "f_sharpe.png",  "Excess return per unit of total volatility"),
          ("Sortino", "1994", "f_sortino.png", "Excess return per unit of downside volatility only"),
          ("Treynor", "1965", "f_treynor.png", "Excess return per unit of systematic risk (β)"),
          ("Information", "1989", "f_info.png", "Consistency of outperformance versus benchmark"),
          ("Calmar", "1991", "f_calmar.png",  "Annual return relative to maximum drawdown")]
pos = [(0.9, 2.25), (4.85, 2.25), (8.8, 2.25), (2.85, 4.78), (6.8, 4.78)]
for (name, yr, f, desc), (cx, cyy) in zip(ratios, pos):
    card(s, cx, cyy, 3.65, 2.32, AMBER_T)
    text(s, cx + 0.25, cyy + 0.2, 3.15, 0.35, [{"runs": [
        (name, dict(font=HEAD, size=15.5, bold=True, color=NAVY)),
        ("  ·  " + yr, dict(font=BODY, size=11.5, color=AMBER_D, bold=True))]}])
    img_fit(s, A + f, cx + 0.25, cyy + 0.62, 3.15, 0.95)
    text(s, cx + 0.25, cyy + 1.66, 3.15, 0.6, [{"runs": [(desc, dict(size=11, color=GRAY))], "line": 1.08}])
pageno(s, 3)

# ================= SLIDE 4 — NEW METRIC =================
s = slide()
kicker_title(s, "The new metric", "Score = diff × r × b")
text(s, 0.9, 1.7, 5.9, 0.65, [{"runs": [
    ("We fit an exponential model to the historical price — an OLS regression on log-price:",
     dict(size=14, color=INK))], "line": 1.15}])
card(s, 0.9, 2.42, 5.9, 0.95, NAVY_T)
img_fit(s, A + "f_model.png", 1.3, 2.56, 5.1, 0.67)
comp4 = [("r", "Fit quality", "Pearson correlation between log-price and time (|r| ≤ 1)"),
         ("b", "Growth slope", "Daily growth rate (log-growth rate)"),
         ("diff", "Price deviation", "ExpReg(last) − Price(last): signed distance from the trend")]
cy = 3.7
for sym, lab, desc in comp4:
    circle(s, 0.9, cy + 0.06, 0.62, TEAL, sym, tsize=15 if len(sym) > 1 else 18, italic=True)
    text(s, 1.75, cy, 5.05, 0.3, [{"runs": [(lab, dict(font=HEAD, size=14.5, bold=True, color=NAVY))]}])
    text(s, 1.75, cy + 0.34, 5.05, 0.65, [{"runs": [(desc, dict(size=12, color=GRAY))], "line": 1.1}])
    cy += 1.06
card(s, 7.3, 1.7, 5.13, 3.62, WHITE, line=LINE)
img_fit(s, A + "concept.png", 7.55, 1.9, 4.63, 3.22)
card(s, 7.3, 5.55, 5.13, 1.35, TEAL_T)
text(s, 7.6, 5.78, 4.55, 0.95, [{"runs": [
    ("Price below the trend + positive trend → positive score → ", dict(size=13, color=INK)),
    ("predicted UP", dict(size=13, bold=True, color=TEAL_D)),
    ("  (mean reversion).", dict(size=13, color=INK))], "line": 1.2}])
pageno(s, 4)

# ================= SLIDE 5 — EXPERIMENT DESIGN =================
s = slide()
kicker_title(s, "Method", "Experiment design")
stats = [("100", "exchange-traded funds"), ("180", "configurations scanned"),
         ("9", "market regimes"), ("5 × 4", "training × test windows")]
sx = 0.9
for num, lab in stats:
    text(s, sx, 1.8, 2.78, 0.75, [{"runs": [(num, dict(font=HEAD, size=40, bold=True, color=NAVY))]}])
    text(s, sx, 2.55, 2.78, 0.35, [{"runs": [(lab, dict(size=12, color=GRAY))]}])
    sx += 2.91
steps = [("1", "Training window", "Compute all 6 metrics on each of the 100 funds."),
         ("2", "Test window", "Observe direction and realize returns."),
         ("3", "Evaluation", "Directional accuracy and long/short portfolio return.")]
sx = 0.9
for n, lab, desc in steps:
    card(s, sx, 3.35, 3.5, 1.75, NAVY_T)
    circle(s, sx + 0.25, 3.62, 0.6, NAVY, n, tsize=17)
    text(s, sx + 1.0, 3.62, 2.3, 0.35, [{"runs": [(lab, dict(font=HEAD, size=15, bold=True, color=NAVY))]}])
    text(s, sx + 1.0, 4.02, 2.3, 0.9, [{"runs": [(desc, dict(size=12, color=GRAY))], "line": 1.12}])
    if n != "3":
        text(s, sx + 3.5, 3.98, 0.52, 0.5, [{"runs": [("→", dict(font=BODY, size=24, bold=True, color=TEAL))],
             "align": PP_ALIGN.CENTER}])
    sx += 4.015
text(s, 0.9, 5.5, 11.53, 1.4, [
    {"runs": [("Each metric is converted into a directional forecast: ", dict(size=13.5)),
              ("a positive score predicts an increase, a negative score predicts a decrease",
               dict(size=13.5, bold=True, color=NAVY)),
              (".", dict(size=13.5))], "line": 1.25, "space_after": 6},
    {"runs": [("In addition, a dollar-neutral long/short portfolio is constructed — 10 assets on each side, "
               "weights proportional to |score|, with a strict cap on the short side.",
               dict(size=13.5, color=INK))], "line": 1.25}])
pageno(s, 5)

# ================= SLIDE 6 — GLOBAL RESULTS =================
s = slide()
kicker_title(s, "Results", "A regime-dependent edge")
img_fit(s, A + "regimes.png", 0.6, 1.75, 7.3, 5.3)
card(s, 8.25, 1.95, 4.2, 1.32, TEAL_T)
text(s, 8.55, 2.16, 3.6, 0.95, [
    {"runs": [("The new metric wins", dict(font=HEAD, size=14.5, bold=True, color=TEAL_D))], "space_after": 3},
    {"runs": [("Reversals and volatile markets", dict(size=12.5, color=INK))]}])
card(s, 8.25, 3.47, 4.2, 1.32, AMBER_T)
text(s, 8.55, 3.68, 3.6, 0.95, [
    {"runs": [("Classic ratios win", dict(font=HEAD, size=14.5, bold=True, color=AMBER_D))], "space_after": 3},
    {"runs": [("Persistent trends", dict(size=12.5, color=INK))]}])
card(s, 8.25, 4.99, 4.2, 1.95, NAVY_T)
text(s, 8.55, 5.2, 3.6, 1.6, [
    {"runs": [("Key takeaway", dict(font=HEAD, size=13.5, bold=True, color=NAVY))], "space_after": 4},
    {"runs": [("The pattern is not random — it is the signature of a mean-reversion signal "
               "versus a momentum signal.", dict(size=12, color=INK))], "line": 1.18}])
pageno(s, 6)

# ================= SLIDES 7 & 8 — CASE STUDIES =================
def case_slide(n, kicker, kcolor, title, cases, note=None):
    s = slide()
    kicker_title(s, kicker, title, kcolor=kcolor)
    ch = 4.25 if note else 4.85
    for cx, label, spark, acc1, acc2, c1, c2, body, ret in cases:
        card(s, cx, 1.85, 5.62, ch, TEAL_T if kcolor == TEAL else AMBER_T)
        text(s, cx + 0.32, 2.12, 3.7, 0.55, [{"runs": [(label, dict(font=BODY, size=12.5, bold=True, color=kcolor))], "line": 1.05}])
        img_fit(s, A + spark, cx + 4.0, 2.05, 1.35, 0.62)
        text(s, cx + 0.32, 2.75, 5.0, 0.8, [{"runs": [
            (acc1, dict(font=HEAD, size=38, bold=True, color=c1)),
            ("  vs  ", dict(font=BODY, size=16, color=GRAY)),
            (acc2, dict(font=HEAD, size=38, bold=True, color=c2))]}])
        text(s, cx + 0.32, 3.62, 5.0, 0.3, [{"runs": [
            ("directional accuracy — new metric vs classic ratios", dict(size=10.5, color=GRAY))]}])
        text(s, cx + 0.32, 4.05, 5.0, 1.35, [{"runs": [(body, dict(size=12, color=INK))], "line": 1.18}])
        if ret:
            text(s, cx + 0.32, 5.55, 5.0, 0.4, [{"runs": [
                ("Long/short return:  ", dict(size=12.5, bold=True, color=NAVY)), *ret]}])
    if note:
        card(s, 0.9, 6.42, 11.53, 0.72, NAVY_T)
        text(s, 1.2, 6.62, 10.9, 0.4, [{"runs": [(note, dict(size=12.5, italic=True, color=NAVY))]}])
    pageno(s, n)

case_slide(7, "Case studies · the new metric wins", TEAL, "When reversion pays", [
    (0.9, "CASE A · COVID CRASH REBOUND (2020)", "spark_a.png", "96.9%", "9.2%", TEAL_D, GRAY,
     "The training window ended near the March 2020 bottom. Almost all ETFs sat below their exponential "
     "trend → large positive diff → predicted UP. The classic ratios, depressed by catastrophic returns, "
     "forecast continued decline.",
     [("+10.9%", dict(size=12.5, bold=True, color=TEAL_D)), ("  vs  ", dict(size=11, color=GRAY)),
      ("−8.8%", dict(size=12.5, bold=True, color=GRAY))]),
    (6.81, "CASE B · 2022 CHOPPY BEAR TAPE (Q4)", "spark_b.png", "90.0%", "22.6%", TEAL_D, GRAY,
     "A market with no clear direction — mean-reverting. Extended names below trend snapped back; names "
     "above it faded. The advantage is not unique to a single crash.",
     [("+8.1%", dict(size=12.5, bold=True, color=TEAL_D)), ("  vs  ", dict(size=11, color=GRAY)),
      ("+2.4%", dict(size=12.5, bold=True, color=GRAY))])])

case_slide(8, "Case studies · the classic ratios win", AMBER_D, "When the trend persists", [
    (0.9, "CASE C · POST-COVID MOMENTUM RALLY", "spark_c.png", "2.0%", "93.9%", GRAY, AMBER_D,
     "End of 2020 — the market had risen for months. Most ETFs were above their trend → negative diff → "
     "predicted a pullback that never came. Momentum continued; the classic ratios were right in almost "
     "every case.",
     None),
    (6.81, "CASE D · 2022 SUSTAINED DECLINE", "spark_d.png", "8.1%", "79.0%", GRAY, AMBER_D,
     "A prolonged, directional decline. Names below trend kept falling — the mean-reversion signal "
     "“caught a falling knife.”",
     [("−6.4%", dict(size=12.5, bold=True, color=GRAY)), ("  vs  ", dict(size=11, color=GRAY)),
      ("+7.8%", dict(size=12.5, bold=True, color=AMBER_D))])],
    note="When the trend persists, mean-reversion forecasting is systematically wrong — the mirror image of Case A.")

# ================= SLIDE 9 — DISCUSSION =================
s = slide()
kicker_title(s, "Discussion", "Two complementary signal types")
def trait_card(x, tint, title, tag, tagc, rows):
    card(s, x, 1.95, 5.62, 2.4, tint)
    paras = [
        {"runs": [(title, dict(font=HEAD, size=17, bold=True, color=tagc))], "space_after": 2},
        {"runs": [(tag, dict(size=11, bold=True, color=tagc))], "space_after": 10}]
    for k, v in rows:
        paras.append({"runs": [(k + "   ", dict(size=13, bold=True, color=NAVY)),
                               (v, dict(size=13, color=INK))], "line": 1.2, "space_after": 6})
    text(s, x + 0.32, 2.25, 5.0, 2.0, paras)
trait_card(0.9, TEAL_T, "The new metric", "MEAN REVERSION", TEAL_D, [
    ("Signal", "positive when price sits below a stable trend"),
    ("Predicts", "a return to the trend"),
    ("Wins in", "reversals and volatile markets (Cases A, B)")])
trait_card(6.81, AMBER_T, "The classic ratios", "MOMENTUM", AMBER_D, [
    ("Signal", "large after strong risk-adjusted returns"),
    ("Predicts", "continuation of the move"),
    ("Wins in", "sustained trends (Cases C, D)")])
card(s, 0.9, 4.85, 11.53, 1.75, NAVY_T)
text(s, 1.22, 5.13, 10.9, 1.3, [
    {"runs": [("The literature (Jegadeesh & Titman, 1993; Daniel & Moskowitz, 2016) documents "
               "“momentum crashes” at reversals — Case A is a textbook example.", dict(size=13, color=INK))],
     "line": 1.22, "space_after": 6},
    {"runs": [("Takeaway: ", dict(size=13, bold=True, color=NAVY)),
              ("a regime-aware combination — switching between the two signal types according to volatility, "
               "trend strength, or the size of diff — is expected to outperform either family on its own.",
               dict(size=13, color=NAVY, bold=True))], "line": 1.22}])
pageno(s, 9)

# ================= SLIDE 10 — CONCLUSIONS (dark) =================
s = slide(NAVY_BG)
text(s, 0.9, 0.5, 11.5, 0.3, [{"runs": [("CONCLUSIONS  ·  FUTURE WORK",
     dict(font=BODY, size=12, bold=True, color=RGBColor(0x4F, 0xD1, 0xCC)))]}])
text(s, 0.9, 0.82, 11.5, 0.75, [{"runs": [("Complementary, not competing",
     dict(font=HEAD, size=30, bold=True, color=WHITE))]}])
concl = [
    [("diff × r × b is a mean-reversion signal", dict(size=13.5, bold=True, color=WHITE)),
     (" — not a replacement for the classic ratios.", dict(size=13.5, color=ICE))],
    [("The classic ratios behave as momentum signals", dict(size=13.5, bold=True, color=WHITE)),
     (" — direct empirical evidence.", dict(size=13.5, color=ICE))],
    [("The two families are complementary", dict(size=13.5, bold=True, color=WHITE)),
     (" — each wins where the other fails.", dict(size=13.5, color=ICE))],
    [("Regime-based combination", dict(size=13.5, bold=True, color=WHITE)),
     (" outperforms any single metric on its own.", dict(size=13.5, color=ICE))]]
cy = 2.0
for i, runs in enumerate(concl, 1):
    circle(s, 0.9, cy, 0.52, TEAL, str(i), tsize=15)
    text(s, 1.65, cy + 0.02, 5.3, 1.0, [{"runs": runs, "line": 1.18}], anchor=MSO_ANCHOR.TOP)
    cy += 1.18
card(s, 7.55, 2.0, 4.9, 4.55, PANEL, radius=0.045)
text(s, 7.95, 2.32, 4.1, 4.0, [
    {"runs": [("NEXT STEP", dict(size=11, bold=True, color=RGBColor(0x4F, 0xD1, 0xCC)))], "space_after": 2},
    {"runs": [("Walk-forward backtest", dict(font=HEAD, size=18, bold=True, color=WHITE))], "space_after": 8},
    {"runs": [("So far we tested each period separately, like snapshots. The next step is to run the test "
               "the way a real investor would: move forward one day at a time, recalculate the metrics, "
               "and update the portfolio as we go. This gives one continuous performance line over the years.",
               dict(size=12.5, color=ICE))], "line": 1.25, "space_after": 10},
    {"runs": [("It will show two things:", dict(size=12.5, color=ICE))], "line": 1.2, "space_after": 6},
    {"runs": [("·  ", dict(size=12.5, color=RGBColor(0x4F, 0xD1, 0xCC), bold=True)),
              ("does our edge survive real, ongoing trading", dict(size=12.5, color=ICE))],
     "line": 1.18, "space_after": 6},
    {"runs": [("·  ", dict(size=12.5, color=RGBColor(0x4F, 0xD1, 0xCC), bold=True)),
              ("does the profit hold up after trading costs", dict(size=12.5, color=ICE))],
     "line": 1.18}])
text(s, 0.9, 6.85, 11.5, 0.35, [{"runs": [
    ("Nehoray Sade · Aviv Asraf   |   Supervised by Dr. Dvir Ross   |   Sami Shamoon College of Engineering",
     dict(size=10.5, color=ICE))]}])

prs.save("Risk-Adjusted-Performance-Analysis—Styled.pptx")
print("saved", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
